#!/usr/bin/env python3
"""Genera Exposicion_BD_GUIOSPRO.pptx con diagramas profesionales."""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("Instale: pip install python-pptx")
    raise SystemExit(1)

ROOT = Path(__file__).resolve().parent
DIAG = ROOT / "diagramas"
OUT = ROOT / "Exposicion_BD_GUIOSPRO.pptx"

TEAL = RGBColor(15, 118, 110)
DARK = RGBColor(15, 23, 42)
GRAY = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
MINT = RGBColor(204, 251, 241)


def svg_to_png(svg_path: Path) -> Path | None:
    png_path = svg_path.with_suffix(".png")
    try:
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPM

        drawing = svg2rlg(str(svg_path))
        if drawing:
            renderPM.drawToFile(drawing, str(png_path), fmt="PNG")
            return png_path if png_path.exists() else None
    except Exception:
        pass
    return None


def ensure_pngs():
    for svg in DIAG.glob("*.svg"):
        png = svg.with_suffix(".png")
        if not png.exists():
            svg_to_png(svg)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def bar_accent(slide):
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.1))
    s.fill.solid()
    s.fill.fore_color.rgb = TEAL
    s.line.fill.background()


def title_slide(prs, title, sub=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, TEAL)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(8.6), Inches(2.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(20)
        p2.font.color.rgb = MINT
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(14)


def section_slide(prs, num, title, sub=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    nb = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(2), Inches(1))
    np = nb.text_frame.paragraphs[0]
    np.text = num
    np.font.size = Pt(72)
    np.font.bold = True
    np.font.color.rgb = RGBColor(45, 212, 191)
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(8.6), Inches(1.5))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(36)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    if sub:
        tp2 = tb.text_frame.add_paragraph()
        tp2.text = sub
        tp2.font.size = Pt(18)
        tp2.font.color.rgb = MINT


def graphic_slide(prs, title, subtitle, image_name, bullets=None, note="", dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK if dark else WHITE)
    if not dark:
        bar_accent(slide)

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(9), Inches(0.55))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = MINT if dark else TEAL

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.55), Inches(0.85), Inches(9), Inches(0.35))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = GRAY if not dark else RGBColor(148, 163, 184)

    png = DIAG / f"{image_name}.png"
    svg = DIAG / f"{image_name}.svg"
    if not png.exists() and svg.exists():
        svg_to_png(svg)

    top = 1.25
    if bullets:
        bb = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(3.2), Inches(5))
        tf = bb.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = b
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE if dark else DARK
            p.space_after = Pt(8)
        img_left = Inches(3.85)
        img_w = Inches(5.9)
    else:
        img_left = Inches(0.55)
        img_w = Inches(9)

    if png.exists():
        slide.shapes.add_picture(str(png), img_left, Inches(top), width=img_w)

    if note:
        nb = slide.shapes.add_textbox(Inches(0.55), Inches(6.85), Inches(9), Inches(0.45))
        np = nb.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(11)
        np.font.italic = True
        np.font.color.rgb = GRAY


def build():
    ensure_pngs()
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(prs, "GUIOSPRO_FLOSS", "Diseño y uso de la base de datos\nMétodo GUIOSAD · Exposición")

    graphic_slide(
        prs,
        "1. Introducción",
        "De archivos CSV a un modelo relacional persistente",
        "csv_vs_bd",
        bullets=[
            "GUIOSPRO evalúa adopción de software FLOSS (método GUIOSAD).",
            "Problema: solo CSV — sin usuarios, historial ni re-evaluaciones.",
            "Objetivo: explicar diseño y uso de la base de datos.",
        ],
    )

    section_slide(prs, "02", "Modelo de datos", "Nivel conceptual")

    graphic_slide(
        prs,
        "Entidades del negocio",
        "Organización · Usuario · Catálogo · Evaluación · Respuestas · Auditoría",
        "er_conceptual",
        note="Catálogo fijo (GUIOSAD) + evaluaciones variables por cliente.",
    )

    graphic_slide(
        prs,
        "Diseño lógico — 9 tablas agrupadas",
        "Seguridad · Catálogo · Evaluaciones · Trazabilidad",
        "nueve_tablas",
        dark=True,
    )

    graphic_slide(
        prs,
        "Flujo de datos — Evaluación de Odoo",
        "Del borrador a la recomendación A / B / C",
        "flujo_odoo",
    )

    graphic_slide(
        prs,
        "Tres reglas de negocio en la BD",
        "Integridad referencial y trazabilidad",
        "reglas_negocio",
        note="Ejemplo: Eval #10 (Odoo) → Eval #25 (padre_id=10, reevaluación #2).",
    )

    title_slide(prs, "Gracias", "GUIOSPRO_FLOSS · SQLite · SQLAlchemy")

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print("Generando PNG desde SVG...")
    ensure_pngs()
    print(f"Presentación: {build()}")
